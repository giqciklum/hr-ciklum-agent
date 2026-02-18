# build_index.py (Versión v25 - Migrated to Google Gemini)
from __future__ import annotations
import os
import re
import glob
import json
import shutil
import base64
import logging
import time
import fitz  # PyMuPDF
import pptx
import docx
from concurrent.futures import ProcessPoolExecutor, as_completed
from typing import List
from tqdm import tqdm
from dotenv import load_dotenv

from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_core.documents import Document
from langchain_core.messages import HumanMessage, SystemMessage
from langchain_chroma import Chroma
from langchain.text_splitter import RecursiveCharacterTextSplitter

# --- Configuración Global y Logging ---
load_dotenv()
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

# --- Constantes y Modelos ---
VISION_MODEL = "gemini-2.0-flash"
EMBEDDING_MODEL = "models/gemini-embedding-001"
DOCS_FOLDER = "docs"
PERSIST_DIR = "chroma_db_v2"
VISION_TIMEOUT = 120
MAX_WORKERS = 1
API_KEY = os.getenv("GOOGLE_API_KEY")

if not API_KEY:
    logging.critical("FATAL: No se ha encontrado la 'GOOGLE_API_KEY'.")
    exit()

# --- Modelos y Splitter ---
vision_llm = ChatGoogleGenerativeAI(
    model=VISION_MODEL,
    google_api_key=API_KEY,
    max_output_tokens=4096,
    convert_system_message_to_human=True,
)
embedder = HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")
text_splitter = RecursiveCharacterTextSplitter(chunk_size=1500, chunk_overlap=200, separators=["\n\n### ", "\n\n## ", "\n\n# ", "\n\n", "\n", " "])


# ==============================================================================
# === LÓGICA DE ENRIQUECIMIENTO ESTRUCTURAL ===
# ==============================================================================

TOPIC_RULES = {
    "EXAMEN_DE_SALUD": {
        "title": "\n\n### PROCESO ESPECÍFICO: EXAMEN DE SALUD (VOLUNTARIO)\n\n",
        "keywords": ["examen médico", "examen de salud", "revisión médica", "aceptar o rechazar el examen", "reconocimiento médico"],
    },
    "PRL_FORMACION": {
        "title": "\n\n### PROCESO ESPECÍFICO: FORMACIÓN EN PREVENCIÓN DE RIESGOS LABORALES (PRL)\n\n",
        "keywords": ["prevención de riesgos laborales", "preventiam", "accesoaula.com", "formación de 120 minutos"],
    },
    "SEGURO_MEDICO": {
        "title": "\n\n### BENEFICIO CLAVE: ALTA EN SEGURO MÉDICO PRIVADO (MAPFRE)\n\n",
        "keywords": ["seguro médico", "mapfre", "alta en el seguro", "formulario de google para el seguro"],
    },
    "RETRIBUCION_FLEXIBLE": {
        "title": "\n\n### BENEFICIO CLAVE: PLAN DE RETRIBUCIÓN FLEXIBLE (EDENRED)\n\n",
        "keywords": ["retribución flexible", "edenred", "tarjeta restaurante", "tarjeta transporte", "ticket restaurant", "guardería"],
    },
    "VACACIONES_Y_PERMISOS": {
        "title": "\n\n### POLÍTICA CLAVE: VACACIONES Y PERMISOS\n\n",
        "keywords": ["vacaciones", "días de vacaciones", "extra agreement days", "solicitar vacaciones", "sesame planner"],
    },
    "BAJA_LABORAL": {
        "title": "\n\n### POLÍTICA CLAVE: BAJA LABORAL (INCAPACIDAD TEMPORAL)\n\n",
        "keywords": ["baja laboral", "baja médica", "incapacidad temporal", "documented sick leave", "undocumented sick leave", "parte de baja"],
    }
}

def enrich_text_with_structural_headings(text: str) -> str:
    enriched_paragraphs = []
    paragraphs = re.split(r'\n\s*\n', text)
    last_inserted_topic = None

    for p in paragraphs:
        if not p.strip():
            continue
        p_lower = p.lower()
        matched_topic = None

        for topic, rules in TOPIC_RULES.items():
            if any(kw in p_lower for kw in rules["keywords"]):
                matched_topic = topic
                break

        if matched_topic and matched_topic != last_inserted_topic:
            enriched_paragraphs.append(TOPIC_RULES[matched_topic]["title"].strip())
            last_inserted_topic = matched_topic

        enriched_paragraphs.append(p)

    return "\n\n".join(enriched_paragraphs)


# --- Utilidades ---
def vision_extract(img_bytes: bytes) -> str:
    try:
        b64_image = base64.b64encode(img_bytes).decode('utf-8')
        resp = vision_llm.invoke([
            HumanMessage(content=[
                {"type": "text", "text": "Eres un experto en OCR. Extrae todo el texto y la estructura de tablas de esta imagen, formateando las tablas en Markdown limpio. Mantén la estructura original lo mejor posible."},
                {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{b64_image}"}}
            ])
        ])
        return enrich_text_with_structural_headings(resp.content.strip())
    except Exception as e:
        logging.error(f"Error en la API de Visión: {e}")
        return ""

def extract_document_title(basename: str) -> str:
    title = os.path.splitext(basename)[0]
    title = re.sub(r'[\W_]+', ' ', title)
    return re.sub(r'\s+', ' ', title).strip().title()

def add_documents_to_vectorstore(docs: List[Document], vector_store: Chroma):
    if not docs: return
    batch_size = 5
    for i in range(0, len(docs), batch_size):
        batch = docs[i:i+batch_size]
        texts = [doc.page_content for doc in batch]
        metadatas = [doc.metadata for doc in batch]
        try:
            vector_store.add_texts(texts=texts, metadatas=metadatas)
        except Exception as e:
            logging.warning(f"Batch failed, retrying in 10s: {e}")
            time.sleep(10)
            try:
                vector_store.add_texts(texts=texts, metadatas=metadatas)
            except Exception as e2:
                logging.error(f"Batch failed permanently: {e2}")
        time.sleep(2)

def has_complex_layout(text: str, line_threshold: int = 15, short_line_chars: int = 60) -> bool:
    lines = text.splitlines()
    if not lines: return False
    avg_line_len = sum(len(line) for line in lines) / len(lines) if len(lines) > 0 else 0
    return len(lines) > line_threshold and avg_line_len < short_line_chars

# --- Procesadores de Ficheros ---

def process_pdf(path: str) -> List[Document]:
    basename = os.path.basename(path)
    doc_title = extract_document_title(basename)
    documents = []

    try:
        doc = fitz.open(path)
        for i, page in enumerate(tqdm(doc, desc=f"PDF: {basename}", leave=False), 1):
            meta = {"source": basename, "page": i, "title": doc_title}
            page_text = page.get_text("text").strip()

            if not page_text or has_complex_layout(page_text):
                logging.info(f"Página {i} de '{basename}' es compleja/imagen. Usando OCR de visión.")
                try:
                    pix = page.get_pixmap(dpi=200)
                    img_bytes = pix.tobytes("png")
                    vision_text = vision_extract(img_bytes)
                    if vision_text:
                        content = f"# Documento: {doc_title}\n## Página: {i}\n\n{vision_text}"
                        documents.append(Document(page_content=content, metadata=meta))
                except Exception as e:
                    logging.warning(f"Error procesando página {i} con OCR: {e}")
            else:
                enriched_text = enrich_text_with_structural_headings(page_text)
                content = f"# Documento: {doc_title}\n## Página: {i}\n\n{enriched_text}"
                documents.append(Document(page_content=content, metadata=meta))
    except Exception as e:
        logging.error(f"Error crítico procesando PDF '{basename}': {e}")
    return documents

def docx_extractor(path: str) -> str:
    doc = docx.Document(path)
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        table_md = "\n".join([f"| {' | '.join(cell.text.strip() for cell in row.cells)} |" for row in table.rows])
        parts.append(f"\n--- TABLA ---\n{table_md}\n--- FIN TABLA ---\n")
    raw_text = "\n\n".join(parts)
    return enrich_text_with_structural_headings(raw_text)

def pptx_extractor(path: str) -> str:
    prs = pptx.Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text)
        enriched_slide_text = enrich_text_with_structural_headings("\n".join(slide_texts))
        parts.append(f"\n## Diapositiva {i}\n{enriched_slide_text}")
    return "\n".join(parts)

def process_generic(path: str, extractor: callable) -> List[Document]:
    basename = os.path.basename(path)
    doc_title = extract_document_title(basename)
    documents = []
    try:
        full_text = extractor(path)
        if full_text:
            meta = {"source": basename, "title": doc_title}
            content = f"# Documento: {doc_title}\n\n{full_text}"
            documents.append(Document(page_content=content, metadata=meta))
    except Exception as e:
        logging.error(f"Error procesando '{basename}': {e}")
    return documents

def process_file(path: str) -> List[Document]:
    ext = path.lower().split('.')[-1]
    processor_map = {
        "pdf": process_pdf,
        "docx": lambda p: process_generic(p, docx_extractor),
        "pptx": lambda p: process_generic(p, pptx_extractor),
    }
    if ext in processor_map:
        return processor_map[ext](path)
    logging.warning(f"Extensión '{ext}' no soportada para el fichero {os.path.basename(path)}")
    return []

# --- Main Execution ---
if __name__ == "__main__":
    print(f"── Construyendo índice vectorial v25 (Gemini + Enriquecido) en '{PERSIST_DIR}' ──")
    if os.path.exists(PERSIST_DIR):
        print(f"Borrando el directorio de índice antiguo: {PERSIST_DIR}")
        shutil.rmtree(PERSIST_DIR)

    db = Chroma(persist_directory=PERSIST_DIR, embedding_function=embedder)

    extensions = ("*.pdf", "*.pptx", "*.docx")
    files = [f for ext in extensions for f in glob.glob(os.path.join(DOCS_FOLDER, f"**/{ext}"), recursive=True)]

    print(f"Se encontraron {len(files)} documentos. Procesando y enriqueciendo automáticamente...")
    total_chunks = 0
    with ProcessPoolExecutor(max_workers=MAX_WORKERS) as pool:
        futs = {pool.submit(process_file, fp): fp for fp in files}
        for fut in tqdm(as_completed(futs), total=len(futs), desc="Procesando Ficheros"):
            try:
                raw_documents = fut.result()
                if raw_documents:
                    chunks = text_splitter.split_documents(raw_documents)
                    add_documents_to_vectorstore(chunks, db)
                    total_chunks += len(chunks)
            except Exception as e:
                logging.error(f"❌ Error crítico en el futuro del fichero {futs[fut]}: {e}", exc_info=True)

    print(f"\nSe generaron e indexaron un total de {total_chunks} chunks semánticos.")
    print(f"✅ Índice vectorial enriquecido creado con éxito y guardado en '{PERSIST_DIR}'.")
